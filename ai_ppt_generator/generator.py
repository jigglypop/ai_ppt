import os
import argparse
from pathlib import Path
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# .env 파일에서 환경 변수 로드
load_dotenv()

class AIPresentation:
    def __init__(self, api_key=None):
        # 제공된 API 키 사용 또는 환경 변수에서 가져오기
        self.api_key = api_key or os.getenv("OPENAI_API_KEY")
        if not self.api_key:
            raise ValueError("OpenAI API 키가 필요합니다. OPENAI_API_KEY 환경 변수를 설정하거나 매개변수로 전달하세요.")
        
        self.client = OpenAI(api_key=self.api_key)
    
    def generate_outline(self, topic, num_slides=5):
        """GPT를 사용하여 프레젠테이션 개요 생성"""
        prompt = f"""
        "{topic}"에 관한 파워포인트 프레젠테이션의 상세 개요를 작성해주세요.
        개요는 다음을 포함한 {num_slides}개의 슬라이드로 구성되어야 합니다:
        1. 제목 슬라이드
        2. 소개
        3-{num_slides-1}. 내용 슬라이드 (글머리 기호 포함)
        {num_slides}. 결론/요약 슬라이드
        
        다음 구조의 JSON 객체 형식으로 응답해주세요:
        {{
            "title": "프레젠테이션 제목",
            "slides": [
                {{
                    "title": "슬라이드 1 제목",
                    "content": ["글머리 기호 1", "글머리 기호 2", ...],
                    "notes": "이 슬라이드에 대한 발표자 노트"
                }},
                ...
            ]
        }}
        """
        
        response = self.client.chat.completions.create(
            model="gpt-4o",
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": "당신은 전문적인 프레젠테이션 제작자입니다. 상세하고 유익한 프레젠테이션을 만드세요."},
                {"role": "user", "content": prompt}
            ]
        )
        
        return response.choices[0].message.content
    
    def create_presentation(self, outline_json, output_path="presentation.pptx"):
        """생성된 개요에서 파워포인트 프레젠테이션 생성"""
        import json
        outline = json.loads(outline_json)
        
        # 프레젠테이션 생성
        prs = Presentation()
        
        # 제목 슬라이드
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = outline["title"]
        subtitle.text = "AI로 생성됨"
        
        # 내용 슬라이드
        for slide_data in outline["slides"]:
            content_slide_layout = prs.slide_layouts[1]  # 제목 및 내용 레이아웃
            slide = prs.slides.add_slide(content_slide_layout)
            
            # 제목 설정
            title = slide.shapes.title
            title.text = slide_data["title"]
            
            # 글머리 기호 추가
            content = slide.placeholders[1]
            text_frame = content.text_frame
            
            for idx, bullet_point in enumerate(slide_data["content"]):
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet_point
                p.level = 0
            
            # 발표자 노트 추가 (가능한 경우)
            if "notes" in slide_data and slide_data["notes"]:
                slide.notes_slide.notes_text_frame.text = slide_data["notes"]
        
        # 프레젠테이션 저장
        prs.save(output_path)
        return output_path 